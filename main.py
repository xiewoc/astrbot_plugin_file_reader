from astrbot.api.event import filter, AstrMessageEvent          # pyright: ignore[reportMissingImports]
from astrbot.api.star import Context, Star, register            # pyright: ignore[reportMissingImports]
from astrbot.api.provider import ProviderRequest                # pyright: ignore[reportMissingImports]
import astrbot.api.message_components as Comp                   # pyright: ignore[reportMissingImports] 
from astrbot.api.all import *                                   # pyright: ignore[reportMissingImports] 
from astrbot.api import logger                                  # pyright: ignore[reportMissingImports]

import os
from datetime import datetime as dt  # 正确导入datetime类
from enum import Enum
from dataclasses import dataclass
from typing import Dict, Optional, Callable, Union, Tuple, Set
from pathlib import Path
import tempfile
import traceback
from functools import lru_cache
import warnings

# 可选依赖导入，提供更好的错误信息
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
    PDFMINER_AVAILABLE = True
except ImportError:
    PDFMINER_AVAILABLE = False
    pdf_extract_text = None

try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except ImportError:
    DOCX2TXT_AVAILABLE = False
    docx2txt = None

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    pd = None

try:
    from docx import Document
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False
    Document = None

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    Presentation = None

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False
    chardet = None

# 常量定义
class FileCategory(Enum):
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    CODE = "code"
    MARKUP = "markup"
    CONFIG = "config"
    TEXT = "text"
    OTHER = "other"

@dataclass
class FileTypeInfo:
    """文件类型信息"""
    extensions: Set[str]
    category: FileCategory
    handler: str
    mime_types: Optional[Set[str]] = None
    requires: Optional[Set[str]] = None  # 需要的库
    max_size: Optional[int] = None  # 最大文件大小（字节）

# 文件类型配置
FILE_TYPES_CONFIG = {
    # 文档格式
    "pdf": FileTypeInfo(
        extensions={".pdf"},
        category=FileCategory.DOCUMENT,
        handler="read_pdf_to_text",
        mime_types={"application/pdf"},
        requires={"pdfminer"}
    ),
    "docx": FileTypeInfo(
        extensions={".docx", ".doc"},
        category=FileCategory.DOCUMENT,
        handler="read_docx_to_text",
        mime_types={
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword"
        },
        requires={"docx2txt", "python-docx"}
    ),
    
    # 电子表格
    "excel": FileTypeInfo(
        extensions={".xlsx", ".xls", ".xlsm", ".ods"},
        category=FileCategory.SPREADSHEET,
        handler="read_excel_to_text",
        mime_types={
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel"
        },
        requires={"pandas"}
    ),
    "csv": FileTypeInfo(
        extensions={".csv"},
        category=FileCategory.SPREADSHEET,
        handler="read_csv_to_text",
        mime_types={"text/csv"},
        requires={"pandas"}
    ),
    
    # 演示文稿
    "pptx": FileTypeInfo(
        extensions={".pptx", ".ppt", ".odp"},
        category=FileCategory.PRESENTATION,
        handler="read_pptx_to_text",
        mime_types={
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint"
        },
        requires={"python-pptx"}
    ),
    
    # 文本文件（包括代码、配置等）
    "text": FileTypeInfo(
        extensions={
            ".txt", ".log", ".md", ".markdown", ".json", ".xml", ".html", ".htm",
            ".yaml", ".yml", ".ini", ".cfg", ".conf", ".properties", ".env",
            ".py", ".java", ".cpp", ".c", ".h", ".hpp", ".cs", ".js", ".ts",
            ".php", ".rb", ".go", ".rs", ".swift", ".kt", ".scala", ".sh",
            ".bash", ".ps1", ".bat", ".cmd", ".vbs", ".sql", ".toml"
        },
        category=FileCategory.TEXT,
        handler="read_text_file",
        mime_types={"text/plain", "text/x-python", "application/json", "text/xml", "text/html"},
        max_size=10 * 1024 * 1024  # 10MB限制
    ),
}

class FileReaderError(Exception):
    """文件读取错误基类"""
    pass

class UnsupportedFileError(FileReaderError):
    """不支持的文件格式"""
    pass

class FileTooLargeError(FileReaderError):
    """文件过大"""
    pass

class DependencyMissingError(FileReaderError):
    """依赖缺失"""
    pass

class FileReader:
    """文件读取器"""
    
    # 扩展名到文件类型的映射（自动生成）
    EXTENSION_TO_TYPE: Dict[str, str] = {}
    for file_type, info in FILE_TYPES_CONFIG.items():
        for ext in info.extensions:
            EXTENSION_TO_TYPE[ext.lower()] = file_type
    
    def __init__(self, max_file_size: Optional[int] = None):
        self.max_file_size = max_file_size or 50 * 1024 * 1024  # 默认50MB
        self._temp_files = []  # 临时文件跟踪
        
    def __del__(self):
        """清理临时文件"""
        self.cleanup_temp_files()
    
    def cleanup_temp_files(self):
        """清理所有临时文件"""
        for temp_file in self._temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except:
                pass
        self._temp_files.clear()
    
    @lru_cache(maxsize=100)
    def detect_file_type(self, file_path: Union[str, Path]) -> Optional[str]:
        """
        检测文件类型
        
        优先级：
        1. 通过python-magic检测MIME类型
        2. 通过文件扩展名
        3. 通过文件内容分析（如果前两者失败）
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
            
        # 方法1：使用python-magic检测MIME类型
        try:
            import magic
            mime_type = magic.from_file(str(file_path), mime=True)
            
            # 查找匹配的MIME类型
            for file_type, info in FILE_TYPES_CONFIG.items():
                if info.mime_types and mime_type in info.mime_types:
                    return file_type
            
            # 特殊处理Office文档
            if "vnd.openxmlformats-officedocument" in mime_type:
                if 'wordprocessingml' in mime_type:
                    return "docx"
                elif 'spreadsheetml' in mime_type:
                    return "excel"
                elif 'presentationml' in mime_type:
                    return "pptx"
                
        except ImportError:
            pass  # python-magic不可用
        
        # 方法2：使用扩展名
        extension = file_path.suffix.lower()
        if extension in self.EXTENSION_TO_TYPE:
            return self.EXTENSION_TO_TYPE[extension]
        
        # 方法3：尝试通过文件内容判断
        try:
            with open(file_path, 'rb') as f:
                header = f.read(1024)  # 读取文件头
            
            # 简单的魔术数字检测
            if header.startswith(b'%PDF'):
                return "pdf"
            elif header.startswith(b'PK\x03\x04'):  # ZIP文件（docx, xlsx, pptx都是ZIP格式）
                # 需要进一步检查内部结构
                return None  # 暂时无法确定具体类型
            elif b'{\\rtf' in header[:100]:
                return "docx"  # RTF文档
            
        except Exception:
            pass
            
        return None
    
    def check_dependencies(self, file_type: str) -> None:
        """检查文件类型所需的依赖是否可用"""
        info = FILE_TYPES_CONFIG.get(file_type)
        if not info or not info.requires:
            return
            
        missing_deps = []
        for dep in info.requires:
            if dep == "pdfminer" and not PDFMINER_AVAILABLE:
                missing_deps.append("pdfminer.six")
            elif dep == "docx2txt" and not DOCX2TXT_AVAILABLE:
                missing_deps.append("docx2txt")
            elif dep == "pandas" and not PANDAS_AVAILABLE:
                missing_deps.append("pandas")
            elif dep == "python-docx" and not PYTHON_DOCX_AVAILABLE:
                missing_deps.append("python-docx")
            elif dep == "python-pptx" and not PYTHON_PPTX_AVAILABLE:
                missing_deps.append("python-pptx")
                
        if missing_deps:
            raise DependencyMissingError(
                f"读取{file_type}文件需要安装: {', '.join(missing_deps)}"
            )
    
    def validate_file(self, file_path: Union[str, Path]) -> Tuple[Path, str]:
        """验证文件并返回标准化路径和文件类型"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
            
        if not file_path.is_file():
            raise FileReaderError(f"不是文件: {file_path}")
            
        # 检查文件大小
        file_size = file_path.stat().st_size
        if file_size > self.max_file_size:
            raise FileTooLargeError(
                f"文件过大: {file_size:,} 字节 > {self.max_file_size:,} 字节限制"
            )
            
        # 检测文件类型
        file_type = self.detect_file_type(file_path)
        if not file_type:
            raise UnsupportedFileError(f"不支持的文件格式: {file_path}")
            
        # 检查文件类型特定的大小限制
        info = FILE_TYPES_CONFIG.get(file_type)
        if info and info.max_size and file_size > info.max_size:
            raise FileTooLargeError(
                f"{file_type}文件过大: {file_size:,} 字节 > {info.max_size:,} 字节限制"
            )
            
        return file_path, file_type
    
    def read_text_file(self, file_path: Path, encoding: Optional[str] = None) -> str:
        """读取文本文件，自动检测编码"""
        if not CHARDET_AVAILABLE:
            # 回退到简单的编码检测
            encodings_to_try = ['utf-8', 'gbk', 'gb2312', 'latin1', 'cp1252']
        else:
            encodings_to_try = []
            
        # 如果指定了编码，先尝试它
        if encoding:
            encodings_to_try.insert(0, encoding)
            
        # 添加chardet检测
        if CHARDET_AVAILABLE:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # 只读取部分内容进行检测
                result = chardet.detect(raw_data)
                if result['confidence'] > 0.7:
                    encodings_to_try.insert(0, result['encoding'])
        
        # 默认编码
        if 'utf-8' not in encodings_to_try:
            encodings_to_try.append('utf-8')
            
        # 尝试不同的编码
        last_error = None
        for enc in encodings_to_try:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, LookupError) as e:
                last_error = e
                continue
                
        # 所有编码都失败
        raise FileReaderError(
            f"无法解码文件 {file_path}，尝试的编码: {encodings_to_try}，最后错误: {last_error}"
        )
    
    def read_pdf_to_text(self, file_path: Path) -> str:
        """读取PDF文件"""
        if not PDFMINER_AVAILABLE:
            raise DependencyMissingError("需要安装 pdfminer.six 库")
            
        try:
            # 使用io.StringIO捕获输出
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                text = pdf_extract_text(str(file_path))
                return text.strip()
        except Exception as e:
            raise FileReaderError(f"读取PDF文件失败: {e}")
    
    def convert_doc_to_docx(self, doc_file: Path, docx_file: Path) -> None:
        """转换DOC到DOCX格式"""
        if not PYTHON_DOCX_AVAILABLE:
            raise DependencyMissingError("需要安装 python-docx 库")
            
        try:
            # 注意：python-docx主要处理docx，对doc支持有限
            # 这里只是简单的转换尝试
            doc = Document(str(doc_file))
            doc.save(str(docx_file))
        except Exception as e:
            raise FileReaderError(f"转换DOC到DOCX失败: {e}")
    
    def read_docx_to_text(self, file_path: Path) -> str:
        """读取DOCX/DOC文件"""
        if not DOCX2TXT_AVAILABLE or not PYTHON_DOCX_AVAILABLE:
            raise DependencyMissingError("需要安装 docx2txt 和 python-docx 库")
            
        try:
            file_path_str = str(file_path)
            
            if file_path.suffix.lower() == '.doc':
                # 创建临时文件
                with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                    tmp_docx = tmp.name
                    self._temp_files.append(tmp_docx)
                
                # 转换
                self.convert_doc_to_docx(file_path, Path(tmp_docx))
                
                # 读取转换后的文件
                text = docx2txt.process(tmp_docx)
                
                # 清理临时文件
                try:
                    os.unlink(tmp_docx)
                    if tmp_docx in self._temp_files:
                        self._temp_files.remove(tmp_docx)
                except:
                    pass
                    
            else:
                text = docx2txt.process(file_path_str)
                
            return text.strip()
        except Exception as e:
            raise FileReaderError(f"读取Word文件失败: {e}")
    
    def read_excel_to_text(self, file_path: Path) -> str:
        """读取Excel文件"""
        if not PANDAS_AVAILABLE:
            raise DependencyMissingError("需要安装 pandas 库")
            
        try:
            # 尝试读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            results = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    # 处理空表
                    if df.empty:
                        results.append(f"工作表: {sheet_name} (空)")
                        continue
                    
                    # 转换为文本
                    text = df.to_string(index=False, na_rep='NA')
                    results.append(f"工作表: {sheet_name}\n{text}")
                    
                except Exception as e:
                    results.append(f"工作表: {sheet_name} (读取失败: {str(e)})")
                    continue
            
            return "\n\n" + "="*50 + "\n\n".join(results)
            
        except Exception as e:
            raise FileReaderError(f"读取Excel文件失败: {e}")
    
    def read_csv_to_text(self, file_path: Path) -> str:
        """读取CSV文件"""
        if not PANDAS_AVAILABLE:
            raise DependencyMissingError("需要安装 pandas 库")
            
        try:
            # 尝试自动检测分隔符和编码
            with open(file_path, 'rb') as f:
                sample = f.read(10000).decode('utf-8', errors='ignore')
                
            # 简单的分隔符检测
            if '\t' in sample:
                sep = '\t'
            elif ';' in sample:
                sep = ';'
            else:
                sep = ','
            
            # 读取CSV
            try:
                df = pd.read_csv(file_path, sep=sep, encoding='utf-8')
            except UnicodeDecodeError:
                try:
                    df = pd.read_csv(file_path, sep=sep, encoding='gbk')
                except:
                    df = pd.read_csv(file_path, sep=sep, encoding='latin1')
            
            return df.to_string(index=False)
            
        except Exception as e:
            raise FileReaderError(f"读取CSV文件失败: {e}")
    
    def read_pptx_to_text(self, file_path: Path) -> str:
        """读取PPTX文件"""
        if not PYTHON_PPTX_AVAILABLE:
            raise DependencyMissingError("需要安装 python-pptx 库")
            
        try:
            prs = Presentation(str(file_path))
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                    
                    elif hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                
                if slide_texts:
                    slides_content.append(f"幻灯片 {i}:\n" + "\n".join(slide_texts))
            
            return "\n\n" + "="*50 + "\n\n".join(slides_content)
            
        except Exception as e:
            raise FileReaderError(f"读取PPTX文件失败: {e}")
    
    def read_file(self, file_path: Union[str, Path], 
                  max_length: Optional[int] = None) -> str:
        """
        读取文件的主要入口
        
        Args:
            file_path: 文件路径
            max_length: 最大返回长度（字符数），None表示不限制
            
        Returns:
            文件内容字符串
        """
        self.cleanup_temp_files()  # 清理之前的临时文件
        
        try:
            # 验证文件并获取文件类型
            file_path, file_type = self.validate_file(file_path)
            
            # 检查依赖
            self.check_dependencies(file_type)
            
            # 选择处理函数
            handler_map = {
                "pdf": self.read_pdf_to_text,
                "docx": self.read_docx_to_text,
                "excel": self.read_excel_to_text,
                "csv": self.read_csv_to_text,
                "pptx": self.read_pptx_to_text,
                "text": self.read_text_file,
            }
            
            handler = handler_map.get(file_type)
            if not handler:
                raise UnsupportedFileError(f"没有找到处理 {file_type} 的处理器")
            
            # 读取文件
            content = handler(file_path)
            
            # 限制长度
            if max_length and len(content) > max_length:
                content = content[:max_length] + f"\n\n[内容截断，总长度: {len(content):,} 字符]"
            
            # 添加文件信息
            file_info = (
                f"\n\n{'='*60}\n"
                f"文件名: {file_path.name}\n"
                f"文件类型: {file_type}\n"
                f"文件大小: {file_path.stat().st_size:,} 字节\n"
                f"读取时间: {dt.now().strftime('%Y-%m-%d %H:%M:%S')}\n"  # 使用dt.now()
                f"{'='*60}"
            )
            
            return content + file_info
            
        except FileReaderError:
            raise  # 重新抛出已知错误
        except Exception as e:
            # 记录详细错误信息
            error_msg = f"读取文件时发生未知错误: {e}\n{traceback.format_exc()}"
            raise FileReaderError(error_msg)
        finally:
            # 确保清理临时文件
            self.cleanup_temp_files()

# 全局文件读取器实例
_file_reader = FileReader()

# 向后兼容的函数
def read_any_file_to_text(file_path: str, max_length: Optional[int] = None) -> str:
    """兼容旧版本的函数"""
    try:
        return _file_reader.read_file(file_path, max_length)
    except FileReaderError as e:
        return f"读取文件失败: {str(e)}"
    except Exception as e:
        return f"读取文件时发生错误: {str(e)}"

@register("astrbot_plugin_file_reader", "xiewoc", 
          "一个将文件内容传给llm的插件", "1.0.3", 
          "https://github.com/xiewoc/astrbot_plugin_file_reader")
class astrbot_plugin_file_reader(Star):
    def __init__(self, context: Context):
        super().__init__(context)
        self.current_file_info: Optional[Dict] = None
        self.file_reader = FileReader(max_file_size=100 * 1024 * 1024)  # 限制10MB
        self.reading_history = []  # 读取历史记录
        
    async def _process_file(self, file_item: Comp.File) -> Optional[Dict]:
        """处理单个文件"""
        try:
            # 获取文件
            file_path = await file_item.get_file()
            logger.info(f"接收到文件: {file_path}")
            
            # 读取文件内容
            content = self.file_reader.read_file(file_path, max_length=100000)  # 限制100K字符
            
            # 记录读取历史
            self.reading_history.append({
                "filename": Path(file_path).name,
                "timestamp": dt.now().isoformat(),  # 使用dt.now()
                "size": Path(file_path).stat().st_size if Path(file_path).exists() else 0,
                "success": True
            })
            
            # 限制历史记录长度
            if len(self.reading_history) > 100:
                self.reading_history = self.reading_history[-50:]
            
            return {
                "path": file_path,
                "name": Path(file_path).name,
                "content": content,
                "size": len(content)
            }
            
        except FileNotFoundError:
            logger.error(f"文件不存在: {file_path}")
        except FileTooLargeError as e:
            logger.error(f"文件过大: {str(e)}")
        except UnsupportedFileError as e:
            logger.warning(f"不支持的文件格式: {str(e)}")
        except DependencyMissingError as e:
            logger.error(f"缺少依赖: {str(e)}")
        except Exception as e:
            logger.error(f"读取文件失败: {str(e)}\n{traceback.format_exc()}")
            
        return None
    
    @event_message_type(EventMessageType.ALL)               # type: ignore
    async def on_receive_msg(self, event: AstrMessageEvent):
        """当获取到有文件时"""
        if event.is_at_or_wake_command:
            file_items = []
            
            # 收集所有文件
            for item in event.message_obj.message:
                if isinstance(item, Comp.File):
                    file_items.append(item)
            
            # 处理文件
            for file_item in file_items:
                file_info = await self._process_file(file_item)
                
                if file_info and file_info["content"]:
                    self.current_file_info = file_info
                    logger.info(f"成功读取文件: {file_info['name']}，大小: {len(file_info['content']):,} 字符")
                    
                    # 可以在这里发送确认消息
                    # await event.reply(f"已读取文件: {file_info['name']}")
                else:
                    logger.warning(f"读取文件内容为空或失败")
                    self.current_file_info = None
    
    @filter.on_llm_request()
    async def on_request(self, event: AstrMessageEvent, req: ProviderRequest):
        """将文件内容添加到LLM请求"""
        if self.current_file_info and self.current_file_info["content"]:
            file_info = self.current_file_info
            
            # 格式化文件内容
            file_context = (
                f"\nFile name:{file_info['name']}\n"
                f"\n[SOF]\n"
                f"{file_info['content']}\n"
                f"[EOF]\n"
            )
            
            # 添加到提示词
            req.prompt += file_context
            logger.info(f"已将文件 {file_info['name']} 内容添加到请求中")
            
            # 清空当前文件信息
            self.current_file_info = None
    
    def get_stats(self) -> Dict:
        """获取插件统计信息"""
        return {
            "total_read": len(self.reading_history),
            "success_read": sum(1 for h in self.reading_history if h.get("success", False)),
            "recent_files": self.reading_history[-10:] if self.reading_history else []
        }